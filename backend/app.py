import os
import io
import json
import re
from datetime import datetime, timedelta, date
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import func, case, desc, text as sqltext
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

app = Flask(__name__)
CORS(app)

app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get(
    'DATABASE_URL',
    'postgresql://postgres:postgres@localhost:5432/kara_taskdb'
)
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)

GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '')
GEMINI_MODEL = os.environ.get('GEMINI_MODEL', 'gemini-2.5-flash')
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)


# ---------------- MODEL ----------------
class Task(db.Model):
    __tablename__ = 'tasks'
    id = db.Column(db.Integer, primary_key=True)
    detail = db.Column(db.String(500), nullable=False)
    notes = db.Column(db.Text, nullable=True, default='')
    task_date = db.Column(db.Date, nullable=False)
    status = db.Column(db.String(20), nullable=False, default='Planned')
    created_at = db.Column(db.DateTime, default=datetime.utcnow, nullable=False)
    modified_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow, nullable=False)

    def to_dict(self):
        return {
            'id': self.id,
            'detail': self.detail,
            'notes': self.notes or '',
            'task_date': self.task_date.isoformat() if self.task_date else None,
            'status': self.status,
            'created_at': self.created_at.isoformat() if self.created_at else None,
            'modified_at': self.modified_at.isoformat() if self.modified_at else None,
        }


STATUS_ORDER = ['Planned', 'Completed', 'Cancelled']


def status_sort_expr():
    return case(
        (Task.status == 'Planned', 0),
        (Task.status == 'Completed', 1),
        (Task.status == 'Cancelled', 2),
        else_=3,
    )


# ---------------- ROUTES: TASKS ----------------
@app.route('/api/health')
def health():
    return jsonify({'status': 'ok'})


@app.route('/api/tasks', methods=['GET'])
def list_tasks():
    page = int(request.args.get('page', 1))
    per_page = int(request.args.get('per_page', 10))
    search = request.args.get('search', '').strip()
    search_date = request.args.get('search_date', '').strip()
    sort_by = request.args.get('sort_by', 'default')
    status_filter = request.args.get('status', '').strip()

    query = Task.query
    if search:
        query = query.filter(Task.detail.ilike(f'%{search}%'))
    if search_date:
        try:
            d = datetime.strptime(search_date, '%Y-%m-%d').date()
            query = query.filter(func.date(Task.created_at) == d)
        except ValueError:
            pass
    if status_filter and status_filter in STATUS_ORDER:
        query = query.filter(Task.status == status_filter)

    if sort_by == 'created':
        query = query.order_by(desc(Task.created_at))
    elif sort_by == 'modified':
        query = query.order_by(desc(Task.modified_at))
    elif sort_by == 'status':
        query = query.order_by(status_sort_expr(), desc(Task.created_at))
    else:
        query = query.order_by(status_sort_expr(), desc(Task.created_at))

    total = query.count()
    items = query.offset((page - 1) * per_page).limit(per_page).all()
    return jsonify({
        'items': [t.to_dict() for t in items],
        'total': total, 'page': page, 'per_page': per_page,
        'pages': (total + per_page - 1) // per_page if per_page else 1,
    })


@app.route('/api/tasks/<int:task_id>', methods=['GET'])
def get_task(task_id):
    return jsonify(Task.query.get_or_404(task_id).to_dict())


@app.route('/api/tasks', methods=['POST'])
def create_task():
    data = request.get_json() or {}
    detail = (data.get('detail') or '').strip()
    notes = (data.get('notes') or '').strip()
    task_date_str = data.get('task_date')
    status = data.get('status', 'Planned')

    if not detail or not task_date_str:
        return jsonify({'error': 'detail and task_date are required'}), 400
    if status not in STATUS_ORDER:
        status = 'Planned'
    try:
        task_date = datetime.strptime(task_date_str, '%Y-%m-%d').date()
    except ValueError:
        return jsonify({'error': 'invalid task_date format'}), 400

    t = Task(detail=detail, notes=notes, task_date=task_date, status=status)
    db.session.add(t)
    db.session.commit()
    return jsonify(t.to_dict()), 201


@app.route('/api/tasks/<int:task_id>', methods=['PUT'])
def update_task(task_id):
    t = Task.query.get_or_404(task_id)
    data = request.get_json() or {}

    if 'detail' in data:
        t.detail = (data['detail'] or '').strip() or t.detail
    if 'notes' in data:
        t.notes = (data['notes'] or '').strip()
    if 'task_date' in data and data['task_date']:
        try:
            t.task_date = datetime.strptime(data['task_date'], '%Y-%m-%d').date()
        except ValueError:
            return jsonify({'error': 'invalid task_date format'}), 400
    if 'status' in data and data['status'] in STATUS_ORDER:
        t.status = data['status']

    db.session.commit()
    return jsonify(t.to_dict())


@app.route('/api/tasks/<int:task_id>/status', methods=['PATCH'])
def update_status(task_id):
    t = Task.query.get_or_404(task_id)
    data = request.get_json() or {}
    new_status = data.get('status')
    if new_status not in STATUS_ORDER:
        return jsonify({'error': 'invalid status'}), 400
    t.status = new_status
    db.session.commit()
    return jsonify(t.to_dict())


@app.route('/api/tasks/<int:task_id>', methods=['DELETE'])
def delete_task(task_id):
    t = Task.query.get_or_404(task_id)
    db.session.delete(t)
    db.session.commit()
    return jsonify({'success': True})


# ---------------- DASHBOARD ----------------
@app.route('/api/dashboard/summary', methods=['GET'])
def dashboard_summary():
    total = db.session.query(func.count(Task.id)).scalar() or 0
    counts = dict(db.session.query(Task.status, func.count(Task.id)).group_by(Task.status).all())
    by_status = {s: int(counts.get(s, 0)) for s in STATUS_ORDER}

    today = date.today()
    start = today - timedelta(days=29)
    rows = db.session.query(
        func.date(Task.created_at).label('d'),
        func.count(Task.id)
    ).filter(func.date(Task.created_at) >= start).group_by('d').all()
    rows_map = {r[0].isoformat() if hasattr(r[0], 'isoformat') else str(r[0]): int(r[1]) for r in rows}

    series = []
    for i in range(30):
        day = start + timedelta(days=i)
        series.append({'date': day.isoformat(), 'count': rows_map.get(day.isoformat(), 0)})

    return jsonify({'total': int(total), 'by_status': by_status, 'chart': series})


# ============================================================
#                    PPT GENERATION (DARK BROWN)
# ============================================================
BROWN = {
    'bg':          RGBColor(0xF5, 0xEE, 0xE2),  # warm ivory/cream background
    'card':        RGBColor(0xFF, 0xFB, 0xF3),  # near-white warm
    'card_alt':    RGBColor(0xEC, 0xE2, 0xD0),  # tan
    'espresso':    RGBColor(0x2B, 0x18, 0x10),  # deepest brown, primary text
    'coffee':      RGBColor(0x4A, 0x2C, 0x20),  # rich brown headers
    'mocha':       RGBColor(0x6B, 0x44, 0x23),  # medium brown
    'umber':       RGBColor(0x8B, 0x6F, 0x47),  # muted secondary
    'gold':        RGBColor(0xC9, 0xA9, 0x61),  # accent gold
    'gold_soft':   RGBColor(0xE3, 0xCC, 0x95),  # soft gold tint
    'copper':      RGBColor(0xB8, 0x73, 0x33),  # copper accent
    'parchment':   RGBColor(0xEF, 0xE5, 0xD3),  # subtle separator
    'muted':       RGBColor(0x8C, 0x7B, 0x68),  # muted text
    'white':       RGBColor(0xFF, 0xFF, 0xFF),
}

STATUS_BADGE = {
    'Planned':   {'bg': RGBColor(0xE3, 0xCC, 0x95), 'fg': RGBColor(0x4A, 0x2C, 0x20)},
    'Completed': {'bg': RGBColor(0x6B, 0x44, 0x23), 'fg': RGBColor(0xF5, 0xEE, 0xE2)},
    'Cancelled': {'bg': RGBColor(0xC5, 0xB0, 0x9A), 'fg': RGBColor(0x4A, 0x2C, 0x20)},
}


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color is not None:
        run.font.color.rgb = color
    return tb


def add_section_header(slide, title_text, subtitle_text=''):
    """Top header used on every content slide: thin gold accent bar + title."""
    # gold accent line
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.06), BROWN['gold'])
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.6),
             title_text, size=26, bold=True, color=BROWN['espresso'])
    if subtitle_text:
        add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
                 subtitle_text, size=12, color=BROWN['muted'], italic=True)
    # subtle bottom divider
    add_rect(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.015), BROWN['parchment'])
    # page footer brand
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             'KARA TASK LIST', size=9, bold=True, color=BROWN['umber'])
    add_text(slide, Inches(6.7), Inches(7.05), Inches(6), Inches(0.3),
             'Confidential', size=9, color=BROWN['muted'], align=PP_ALIGN.RIGHT)


def build_ppt(plan, tasks, date_from, date_to, status_filter):
    """Build a professional dark-brown themed PPT (category-based, no task list)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title = plan.get('title', 'Task Report')
    subtitle = plan.get('subtitle', '')
    summary = plan.get('summary', '')
    categories = plan.get('categories', []) or []
    recommendations = plan.get('recommendations', []) or []

    # =========================================================
    # SLIDE 1: COVER
    # =========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BROWN['bg'])

    add_rect(s, Inches(0), Inches(0), Inches(4.5), Inches(7.5), BROWN['espresso'])
    add_rect(s, Inches(4.5), Inches(0), Inches(0.08), Inches(7.5), BROWN['gold'])

    add_rect(s, Inches(0.6), Inches(0.7), Inches(0.4), Inches(0.05), BROWN['gold'])
    add_text(s, Inches(0.6), Inches(0.85), Inches(3.5), Inches(0.4),
             'KARA TASK LIST', size=12, bold=True, color=BROWN['gold'])
    add_text(s, Inches(0.6), Inches(1.2), Inches(3.5), Inches(0.4),
             'Executive Report', size=10, color=BROWN['gold_soft'], italic=True)

    add_text(s, Inches(0.6), Inches(6.4), Inches(3.5), Inches(0.3),
             'PERIOD', size=9, bold=True, color=BROWN['gold'])
    add_text(s, Inches(0.6), Inches(6.7), Inches(3.5), Inches(0.4),
             f'{date_from}  \u2014  {date_to}', size=12, color=BROWN['bg'])

    add_text(s, Inches(5.2), Inches(2.6), Inches(7.8), Inches(0.4),
             'TASK PERFORMANCE REPORT', size=11, bold=True, color=BROWN['copper'])
    add_rect(s, Inches(5.2), Inches(3.05), Inches(0.8), Inches(0.04), BROWN['gold'])
    add_text(s, Inches(5.2), Inches(3.2), Inches(7.8), Inches(1.6),
             title, size=40, bold=True, color=BROWN['espresso'])
    if subtitle:
        add_text(s, Inches(5.2), Inches(4.7), Inches(7.8), Inches(0.8),
                 subtitle, size=15, color=BROWN['mocha'], italic=True)

    if status_filter:
        add_text(s, Inches(5.2), Inches(5.6), Inches(7.8), Inches(0.4),
                 f'Status filter: {status_filter}', size=11, color=BROWN['umber'])
    add_text(s, Inches(5.2), Inches(6.7), Inches(7.8), Inches(0.4),
             f'Prepared {date.today().isoformat()}', size=10, color=BROWN['muted'], italic=True)

    # =========================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # =========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BROWN['bg'])
    add_section_header(s, 'Executive Summary', 'Period-over-period performance at a glance')

    total = len(tasks)
    by_status = {st: 0 for st in STATUS_ORDER}
    for t in tasks:
        by_status[t['status']] = by_status.get(t['status'], 0) + 1
    completion_rate = round((by_status['Completed'] / total) * 100) if total else 0

    kpis = [
        ('TOTAL TASKS', str(total)),
        ('PLANNED', str(by_status.get('Planned', 0))),
        ('COMPLETED', str(by_status.get('Completed', 0))),
        ('COMPLETION', f'{completion_rate}%'),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.4)
    gap = Inches(0.15)
    start_x = Inches(0.6)
    for i, (label, val) in enumerate(kpis):
        left = start_x + (card_w + gap) * i
        add_rect(s, left, Inches(2.0), card_w, card_h, BROWN['card'])
        add_rect(s, left, Inches(2.0), card_w, Inches(0.06), BROWN['gold'])
        add_text(s, left + Inches(0.3), Inches(2.2), card_w - Inches(0.6), Inches(0.35),
                 label, size=9, bold=True, color=BROWN['umber'])
        add_text(s, left + Inches(0.3), Inches(2.55), card_w - Inches(0.6), Inches(0.85),
                 val, size=36, bold=True, color=BROWN['espresso'])

    add_rect(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(2.95), BROWN['card'])
    add_rect(s, Inches(0.6), Inches(3.7), Inches(0.08), Inches(2.95), BROWN['copper'])
    add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.4),
             'SUMMARY', size=10, bold=True, color=BROWN['copper'])
    add_text(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(2.2),
             summary or 'No summary available.', size=13, color=BROWN['espresso'])

    # =========================================================
    # SLIDE 3+: CATEGORY INSIGHTS (dynamic height, space-based pagination)
    # =========================================================
    if categories:
        # ---- layout constants ----
        FONT_SIZE = 10                 # insight body font (pt)
        CARD_W = Inches(12.1)
        CARD_X = Inches(0.6)
        CONTENT_TOP = Inches(2.0)      # first card starts here
        CONTENT_BOTTOM = Inches(6.9)   # cards must not pass this (footer at 7.05)
        GAP_Y = Inches(0.2)
        HEADER_H = Inches(0.8)         # name + badge + divider zone
        BODY_PAD_BOTTOM = Inches(0.2)
        # approx text metrics for a ~11.6" wide box at 10pt
        CHARS_PER_LINE = 150
        LINE_H = Inches(0.19)          # height per wrapped line at 10pt

        def wrapped_line_count(insight_text):
            n_lines = 0
            for raw_line in insight_text.split('\n'):
                if not raw_line.strip():
                    n_lines += 1
                    continue
                n_lines += max(1, -(-len(raw_line) // CHARS_PER_LINE))  # ceil div
            return max(1, n_lines)

        def card_height_for(insight_text):
            body_h = LINE_H * wrapped_line_count(insight_text)
            return HEADER_H + body_h + BODY_PAD_BOTTOM

        # ---- greedy pack categories into pages by available vertical space ----
        pages = []
        current = []
        y = CONTENT_TOP
        for cat in categories:
            insight = cat.get('insight', '') or 'No insight available.'
            ch = card_height_for(insight)
            if current and (y + ch > CONTENT_BOTTOM):
                pages.append(current)
                current = []
                y = CONTENT_TOP
            current.append((cat, ch))
            y += ch + GAP_Y
        if current:
            pages.append(current)

        total_pages = len(pages)
        total_cats = len(categories)

        for pidx, page in enumerate(pages):
            s = prs.slides.add_slide(blank)
            set_slide_bg(s, BROWN['bg'])
            if total_pages > 1:
                subtitle_str = f'Page {pidx + 1} of {total_pages}  \u00b7  {total_cats} categories total'
            else:
                subtitle_str = f'{total_cats} categor{"ies" if total_cats != 1 else "y"} in this period'
            add_section_header(s, 'Category Insights', subtitle_str)

            top = CONTENT_TOP
            for cat, ch in page:
                left = CARD_X
                cat_name = cat.get('name', 'Untitled')
                n = cat.get('task_count', 0)
                insight = cat.get('insight', '') or 'No insight available.'

                add_rect(s, left, top, CARD_W, ch, BROWN['card'])
                add_rect(s, left, top, Inches(0.06), ch, BROWN['copper'])
                add_text(s, left + Inches(0.25), top + Inches(0.15),
                         CARD_W - Inches(1.6), Inches(0.4),
                         cat_name, size=15, bold=True, color=BROWN['espresso'])
                badge_w = Inches(1.15)
                badge_left = left + CARD_W - badge_w - Inches(0.2)
                add_rounded(s, badge_left, top + Inches(0.2), badge_w, Inches(0.32), BROWN['espresso'])
                add_text(s, badge_left, top + Inches(0.22), badge_w, Inches(0.28),
                         f'{n} TASK{"S" if n != 1 else ""}',
                         size=9, bold=True, color=BROWN['gold'], align=PP_ALIGN.CENTER)
                add_rect(s, left + Inches(0.25), top + Inches(0.65),
                         CARD_W - Inches(0.5), Inches(0.02), BROWN['gold'])
                add_text(s, left + Inches(0.25), top + HEADER_H,
                         CARD_W - Inches(0.5), ch - HEADER_H - BODY_PAD_BOTTOM,
                         insight, size=FONT_SIZE, color=BROWN['espresso'])

                top = top + ch + GAP_Y

    # =========================================================
    # SLIDE: RECOMMENDATIONS
    # =========================================================
    if recommendations:
        s = prs.slides.add_slide(blank)
        set_slide_bg(s, BROWN['bg'])
        add_section_header(s, 'Recommendations', 'Actionable next steps for the upcoming period')

        cols = 2
        card_w2 = Inches(5.95)
        card_h2 = Inches(1.3)
        gap_x = Inches(0.2)
        gap_y = Inches(0.2)
        for idx, r in enumerate(recommendations[:6]):
            col = idx % cols
            row = idx // cols
            left = Inches(0.6) + (card_w2 + gap_x) * col
            top = Inches(2.0) + (card_h2 + gap_y) * row
            add_rect(s, left, top, card_w2, card_h2, BROWN['card'])
            add_rect(s, left, top, Inches(0.06), card_h2, BROWN['copper'])
            add_text(s, left + Inches(0.3), top + Inches(0.15), Inches(0.6), Inches(0.4),
                     f'{idx+1:02d}', size=18, bold=True, color=BROWN['copper'])
            add_text(s, left + Inches(0.95), top + Inches(0.2), card_w2 - Inches(1.2), card_h2 - Inches(0.3),
                     r, size=11, color=BROWN['espresso'])

    # =========================================================
    # CLOSING SLIDE
    # =========================================================
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BROWN['espresso'])
    add_rect(s, Inches(6.0), Inches(3.0), Inches(1.3), Inches(0.06), BROWN['gold'])
    add_text(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.0),
             'Thank You', size=54, bold=True, color=BROWN['bg'], align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
             'Questions & Discussion', size=16, color=BROWN['gold_soft'], italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
             'KARA TASK LIST  \u00b7  Generated for executive review', size=9, color=BROWN['umber'], align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def call_gemini_for_plan(tasks, date_from, date_to, status_filter):
    """Ask Gemini to produce an executive briefing grouped by category (no task list)."""

    # --- Group tasks by category (text before " - " in detail) ---
    def extract_category(detail):
        if ' - ' in detail:
            return detail.split(' - ', 1)[0].strip()
        return 'Uncategorized'

    def build_detailed_insight(items):
        lines = [f"{len(items)} task(s) in this category:"]
        for t in items:
            line = f"• [{t['status']}] {t['detail']}"
            n = (t.get('notes') or '').strip()
            if n:
                n_short = n.replace('\n', ' ')
                if len(n_short) > 180:
                    n_short = n_short[:180] + '...'
                line += f" — {n_short}"
            lines.append(line)
        return "\n".join(lines)
    
    grouped = {}
    for t in tasks:
        cat = extract_category(t['detail'])
        grouped.setdefault(cat, []).append(t)

    # --- Fallback if no Gemini API key ---
    if not GEMINI_API_KEY:
        return {
            'title': 'Task Performance Report',
            'subtitle': f'Period {date_from} to {date_to}',
            'summary': f'This report covers {len(tasks)} task(s) across {len(grouped)} categories for the selected period.',
            'categories': [
                {
                    'name': cat,
                    'task_count': len(items),
                    'insight': f'{len(items)} task(s) recorded under {cat}.',
                }
                for cat, items in grouped.items()
            ],
            'recommendations': [
                'Review pending Planned tasks.',
                'Capture lessons learned from Completed items.',
            ],
        }

    # --- Build prompt content: tasks organized per category (context for reasoning only) ---
    category_blocks = []
    for cat, items in grouped.items():
        block_lines = [f"CATEGORY: {cat} ({len(items)} task(s))"]
        for t in items:
            line = f"  - [{t['status']}] {t['task_date']}: {t['detail']}"
            n = (t.get('notes') or '').strip()
            if n:
                n_short = n.replace('\n', ' ').replace('"', "'")
                if len(n_short) > 300:
                    n_short = n_short[:300] + '...'
                line += f"\n    context: {n_short}"
            block_lines.append(line)
        category_blocks.append("\n".join(block_lines))
    tasks_text = "\n\n".join(category_blocks) if category_blocks else "(no tasks)"

    status_text = status_filter if status_filter else 'All'
    category_names = list(grouped.keys())

    system_prompt = (
        "You are a senior management consultant preparing a board-level executive briefing. "
        "Your tone is formal, precise, and outcome-oriented. You synthesize task data into "
        "insights a busy executive can scan in seconds. You never use emojis, exclamation "
        "marks, or casual language. You group work by category to reveal thematic patterns. "
        "Reply with a single valid JSON object and nothing else."
    )

    user_prompt = f"""Produce an executive briefing plan for the following task portfolio, organized by CATEGORY.

PERIOD: {date_from} to {date_to}
STATUS FILTER: {status_text}
TOTAL TASK COUNT: {len(tasks)}
CATEGORIES ({len(category_names)}): {', '.join(category_names)}

TASK PORTFOLIO (grouped by category; each task may include a `context:` line for your reasoning only):
{tasks_text}

Return ONE JSON object with EXACTLY this schema (no extra keys, no comments, no trailing commas):
{{
  "title": "string",
  "subtitle": "string",
  "summary": "string",
  "categories": [
    {{
      "name": "string",
      "task_count": number,
      "insight": "string"
    }}
  ],
  "recommendations": ["string", "string", ...]
}}

FIELD GUIDELINES:

- title: concise, professional, max 7 words. Avoid hype.
- subtitle: one-line scope descriptor, max 14 words.
- summary: 5-7 sentences. Open with the headline finding across the portfolio. Then describe the thematic pattern PER CATEGORY (which category dominates activity, which is stalled, which shows progress). Reference real counts and category names. End with the forward outlook. Formal third-person tone.
- categories: one object per category present in the input, in the same order as the CATEGORIES list above.
    * name: EXACT category name from input (do not rename, do not merge, do not split).
    * task_count: integer count from input.
    * insight: describe the category's overall posture in 1 sentence, THEN provide a per-task breakdown. For EACH task write one line beginning with "• [STATUS] " followed by a short description of the task and, where relevant, any risk/blocker or progress inferred from its notes. Use "\n" to separate lines. Reference the actual task detail.
- recommendations: 3-5 actionable next steps grounded in the data, ideally tied to specific categories. Max 18 words each.

CRITICAL RULES:
- Output ONE valid JSON object. No markdown fences, no prose before or after.
- Per-task detail SHOULD appear inside each category's "insight" field, one line per task.
- Do NOT include a separate "tasks" array. Keep all task detail inside "insight".
- You MAY use notes/context to enrich the per-task description.
- Escape any double quotes inside string values as \\".
- No emojis except the "•" bullet, no filler words.
"""

    def _parse_json_lenient(raw: str):
        """Try to extract a JSON object from Gemini's output even if slightly malformed."""
        s = raw.strip()
        # strip code fences if present
        s = re.sub(r'^```(?:json)?\s*', '', s)
        s = re.sub(r'\s*```$', '', s)
        # first attempt: direct
        try:
            return json.loads(s)
        except json.JSONDecodeError:
            pass
        # second attempt: extract substring from first { to last }
        start = s.find('{')
        end = s.rfind('}')
        if start != -1 and end != -1 and end > start:
            candidate = s[start:end + 1]
            # remove trailing commas before } or ]
            candidate = re.sub(r',(\s*[}\]])', r'\1', candidate)
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                pass
        raise ValueError("Could not parse JSON from Gemini response")

    try:
        try:
            from google.generativeai.types import GenerationConfig
            gen_config_obj = GenerationConfig(
                temperature=0.4,
                max_output_tokens=4096,
                response_mime_type="application/json",
                thinking_config={"thinking_budget": 0},
            )
            model = genai.GenerativeModel(
                model_name=GEMINI_MODEL,
                system_instruction=system_prompt,
                generation_config=gen_config_obj,
            )
        except (TypeError, ImportError):
            gen_config = {
                "temperature": 0.4,
                "max_output_tokens": 4096,
                "response_mime_type": "application/json",
            }
            model = genai.GenerativeModel(
                model_name=GEMINI_MODEL,
                system_instruction=system_prompt,
                generation_config=gen_config,
            )

        resp = model.generate_content(user_prompt)
        raw_text = (resp.text or '').strip()
        result = _parse_json_lenient(raw_text)

        # --- Safety net: enforce category names + counts from source of truth ---
        source_categories = [
            {'name': cat, 'task_count': len(items)}
            for cat, items in grouped.items()
        ]
        gemini_insights = {
            c.get('name'): c.get('insight', '')
            for c in result.get('categories', [])
            if isinstance(c, dict)
        }
        result['categories'] = [
            {
                'name': cat,
                'task_count': len(items),
                'insight': gemini_insights.get(cat) or build_detailed_insight(items),
            }
            for cat, items in grouped.items()
        ]
        # strip any stray "tasks" key if Gemini added it against instructions
        for c in result['categories']:
            c.pop('tasks', None)

        return result

    except Exception as e:
        print(f"[Gemini error] {e}")
        return {
            'title': 'Task Performance Report',
            'subtitle': f'Period {date_from} to {date_to}',
            'summary': f'This report covers {len(tasks)} task(s) across {len(grouped)} categories. (AI generation unavailable: {e})',
            'categories': [
                {
                    'name': cat,
                    'task_count': len(items),
                    'insight': f'{len(items)} task(s) recorded under {cat}.',
                }
                for cat, items in grouped.items()
            ],
            'recommendations': ['Review the task list manually.'],
        }


@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.get_json() or {}
    date_from = data.get('date_from')
    date_to = data.get('date_to')
    status_filter = data.get('status', '').strip()

    if not date_from or not date_to:
        return jsonify({'error': 'date_from and date_to are required'}), 400
    try:
        df = datetime.strptime(date_from, '%Y-%m-%d').date()
        dt = datetime.strptime(date_to, '%Y-%m-%d').date()
    except ValueError:
        return jsonify({'error': 'invalid date format'}), 400

    q = Task.query.filter(Task.task_date >= df, Task.task_date <= dt)
    if status_filter and status_filter in STATUS_ORDER:
        q = q.filter(Task.status == status_filter)
    q = q.order_by(status_sort_expr(), Task.task_date.asc())
    tasks = [t.to_dict() for t in q.all()]

    if not tasks:
        return jsonify({'error': 'No tasks found for the selected range/status.'}), 404

    plan = call_gemini_for_plan(tasks, date_from, date_to, status_filter)
    buf = build_ppt(plan, tasks, date_from, date_to, status_filter)
    filename = f"kara_task_report_{date_from}_to_{date_to}.pptx"
    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )


# ============================================================
#                    EXCEL EXPORT
# ============================================================
def build_excel(tasks, date_from, date_to, status_filter):
    """Build an .xlsx workbook: Tasks sheet (detail rows) + Summary sheet."""
    wb = Workbook()

    # ---- palette (matches the brown theme, hex without leading #) ----
    HEADER_BG = "4A2C20"   # coffee
    HEADER_FG = "F5EEE2"   # ivory
    TITLE_FG = "2B1810"    # espresso
    ACCENT = "C9A961"      # gold
    ZEBRA = "F5EEE2"       # light row tint
    STATUS_FILL = {
        "Planned":   "E3CC95",
        "Completed": "6B4423",
        "Cancelled": "C5B09A",
    }
    STATUS_FONT = {
        "Planned":   "4A2C20",
        "Completed": "F5EEE2",
        "Cancelled": "4A2C20",
    }

    thin = Side(style="thin", color="D8CBB6")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # ================= SHEET 1: TASKS =================
    ws = wb.active
    ws.title = "Tasks"

    # report meta block
    ws["A1"] = "KARA TASK LIST — Task Export"
    ws["A1"].font = Font(name="Calibri", size=15, bold=True, color=TITLE_FG)
    ws["A2"] = f"Period: {date_from} to {date_to}"
    ws["A2"].font = Font(name="Calibri", size=10, color="6B4423")
    ws["A3"] = f"Status filter: {status_filter if status_filter else 'All'}"
    ws["A3"].font = Font(name="Calibri", size=10, color="6B4423")
    ws["A4"] = f"Generated: {date.today().isoformat()}   |   Total: {len(tasks)} task(s)"
    ws["A4"].font = Font(name="Calibri", size=10, color="6B4423")

    header_row = 6
    headers = ["No", "Task Date", "Detail", "Notes", "Status", "Created At", "Modified At"]
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=header_row, column=col, value=h)
        c.font = Font(name="Calibri", size=11, bold=True, color=HEADER_FG)
        c.fill = PatternFill("solid", fgColor=HEADER_BG)
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = border

    def fmt_dt(val):
        if not val:
            return ""
        return str(val).replace("T", " ")[:19]

    r = header_row + 1
    for i, t in enumerate(tasks, start=1):
        row_vals = [
            i,
            t.get("task_date", ""),
            t.get("detail", ""),
            t.get("notes", "") or "",
            t.get("status", ""),
            fmt_dt(t.get("created_at")),
            fmt_dt(t.get("modified_at")),
        ]
        for col, val in enumerate(row_vals, start=1):
            c = ws.cell(row=r, column=col, value=val)
            c.font = Font(name="Calibri", size=10, color=TITLE_FG)
            c.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
            c.border = border
            if i % 2 == 0:
                c.fill = PatternFill("solid", fgColor=ZEBRA)
        # color the status cell
        st = t.get("status", "")
        if st in STATUS_FILL:
            sc = ws.cell(row=r, column=5)
            sc.fill = PatternFill("solid", fgColor=STATUS_FILL[st])
            sc.font = Font(name="Calibri", size=10, bold=True, color=STATUS_FONT[st])
            sc.alignment = Alignment(horizontal="center", vertical="center")
        r += 1

    # column widths
    widths = [5, 14, 46, 50, 13, 20, 20]
    for col, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(col)].width = w

    # freeze header + autofilter over the table
    ws.freeze_panes = ws.cell(row=header_row + 1, column=1)
    last_col = get_column_letter(len(headers))
    ws.auto_filter.ref = f"A{header_row}:{last_col}{max(header_row, r - 1)}"

    # ================= SHEET 2: SUMMARY =================
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = "Summary"
    ws2["A1"].font = Font(name="Calibri", size=15, bold=True, color=TITLE_FG)
    ws2["A2"] = f"Period: {date_from} to {date_to}"
    ws2["A2"].font = Font(name="Calibri", size=10, color="6B4423")

    by_status = {s: 0 for s in STATUS_ORDER}
    for t in tasks:
        by_status[t.get("status")] = by_status.get(t.get("status"), 0) + 1
    total = len(tasks)
    completion = round((by_status["Completed"] / total) * 100) if total else 0

    sum_header = 4
    for col, h in enumerate(["Metric", "Value"], start=1):
        c = ws2.cell(row=sum_header, column=col, value=h)
        c.font = Font(name="Calibri", size=11, bold=True, color=HEADER_FG)
        c.fill = PatternFill("solid", fgColor=HEADER_BG)
        c.border = border

    rows = [
        ("Total Tasks", total),
        ("Planned", by_status.get("Planned", 0)),
        ("Completed", by_status.get("Completed", 0)),
        ("Cancelled", by_status.get("Cancelled", 0)),
        ("Completion Rate", f"{completion}%"),
    ]
    rr = sum_header + 1
    for label, val in rows:
        a = ws2.cell(row=rr, column=1, value=label)
        b = ws2.cell(row=rr, column=2, value=val)
        a.font = Font(name="Calibri", size=10, bold=True, color=TITLE_FG)
        b.font = Font(name="Calibri", size=10, color=TITLE_FG)
        a.border = border
        b.border = border
        rr += 1

    ws2.column_dimensions["A"].width = 22
    ws2.column_dimensions["B"].width = 16

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


@app.route('/api/export-excel', methods=['POST'])
def export_excel():
    data = request.get_json() or {}
    date_from = data.get('date_from')
    date_to = data.get('date_to')
    status_filter = data.get('status', '').strip()

    if not date_from or not date_to:
        return jsonify({'error': 'date_from and date_to are required'}), 400
    try:
        df = datetime.strptime(date_from, '%Y-%m-%d').date()
        dt = datetime.strptime(date_to, '%Y-%m-%d').date()
    except ValueError:
        return jsonify({'error': 'invalid date format'}), 400

    q = Task.query.filter(Task.task_date >= df, Task.task_date <= dt)
    if status_filter and status_filter in STATUS_ORDER:
        q = q.filter(Task.status == status_filter)
    q = q.order_by(status_sort_expr(), Task.task_date.asc())
    tasks = [t.to_dict() for t in q.all()]

    if not tasks:
        return jsonify({'error': 'No tasks found for the selected range/status.'}), 404

    buf = build_excel(tasks, date_from, date_to, status_filter)
    filename = f"kara_task_export_{date_from}_to_{date_to}.xlsx"
    return send_file(
        buf,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    )


# ---------------- INIT ----------------
def init_db():
    with app.app_context():
        db.create_all()
        # Auto-migrate: add notes column if it doesn't exist (for existing DBs)
        try:
            db.session.execute(sqltext("ALTER TABLE tasks ADD COLUMN IF NOT EXISTS notes TEXT DEFAULT ''"))
            db.session.commit()
        except Exception as e:
            db.session.rollback()
            print(f"[init_db] migration skipped: {e}")


if __name__ == '__main__':
    init_db()
    app.run(host='0.0.0.0', port=8888, debug=False)